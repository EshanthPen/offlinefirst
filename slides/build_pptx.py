"""
Build offlinefirst-pitch.pptx from the same content as slides/index.html.

Run with:  python3 slides/build_pptx.py
Outputs:   slides/offlinefirst-pitch.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Brand palette (same as index.html)
INK       = RGBColor(0x20, 0x21, 0x24)
MUTED     = RGBColor(0x5F, 0x61, 0x68)
FAINT     = RGBColor(0x9A, 0xA0, 0xA6)
RULE      = RGBColor(0xDA, 0xDC, 0xE0)
PRIMARY   = RGBColor(0x19, 0x67, 0xD2)
SOFT      = RGBColor(0xE8, 0xF0, 0xFE)
PANEL     = RGBColor(0xF8, 0xF9, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x0A, 0x0A, 0x0A)
OK        = RGBColor(0x1E, 0x8E, 0x3E)
BAD       = RGBColor(0xD9, 0x30, 0x25)

SANS    = "Roboto"
DISPLAY = "Google Sans"      # PowerPoint will fall back to a similar sans
MONO    = "Consolas"

# 16:9 dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, rgb):
    """Solid fill the slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_text(slide, left, top, width, height, text,
             font=SANS, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, tracking=None):
    """Insert a transparent textbox with one paragraph."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if tracking is not None:
        # spacing in 1/100 of a point
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(int(tracking * 100)))
    return tb


def add_multi(slide, left, top, width, height, runs,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.4):
    """Textbox with multiple runs in one paragraph, each run a dict of
    {text, font, size, bold, color}."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for spec in runs:
        if spec.get('break_before'):
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = spec.get('text', '')
        r.font.name = spec.get('font', SANS)
        r.font.size = Pt(spec.get('size', 18))
        r.font.bold = spec.get('bold', False)
        r.font.color.rgb = spec.get('color', INK)
    return tb


def add_rect(slide, left, top, width, height, fill, line=None, radius=None):
    """Add a rectangle (rounded if radius given) with solid fill and no line by default."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    if radius:
        # Adjust the corner radius. The first adjustment governs corner radius
        # as a fraction of the shorter side.
        try:
            adj = s.adjustments
            # Convert EMU radius to fraction of min(width, height)/2
            min_side = min(width, height)
            frac = max(0, min(0.5, radius / min_side))
            adj[0] = frac
        except Exception:
            pass
    s.shadow.inherit = False
    return s


def add_circle(slide, cx, cy, r, fill, line=None, line_width=1.5):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, 2 * r, 2 * r)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_width)
    s.shadow.inherit = False
    return s


def add_line(slide, x1, y1, x2, y2, color=PRIMARY, weight=1.0):
    from pptx.shapes.connector import Connector
    l = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    l.line.color.rgb = color
    l.line.width = Pt(weight)
    return l


def shape_text(shape, text, font=SANS, size=14, bold=False, color=INK,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_notes(slide, header, body, cue):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = header
    r0.font.bold = True
    r0.font.size = Pt(12)

    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    r1.text = body
    r1.font.size = Pt(11)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "CUE: " + cue
    r2.font.size = Pt(11)
    r2.font.italic = True


# ----- slide builders -----

def build_title(slide):
    set_slide_bg(slide, WHITE)
    # Logo mark — rounded primary square + outlined circle to its right
    mark_size = Inches(1.4)
    mark_x = Inches(5.6)
    mark_y = Inches(2.0)
    sq = add_rect(slide, mark_x, mark_y, Inches(0.9), Inches(0.9),
                  PRIMARY, radius=Inches(0.22))
    add_circle(slide, mark_x + Inches(1.25), mark_y + Inches(0.45),
               Inches(0.42), WHITE, line=PRIMARY, line_width=3)
    # Wordmark
    add_text(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(1.6),
             "offlinefirst", font=DISPLAY, size=96, bold=True, color=INK,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    # Tag
    add_text(slide, Inches(0), Inches(4.8), SLIDE_W, Inches(0.6),
             "Education without internet.", font=SANS, size=24, color=MUTED,
             align=PP_ALIGN.CENTER)
    # Credit
    add_text(slide, Inches(0), Inches(6.6), SLIDE_W, Inches(0.4),
             "YCS 2026  ·  OPEN SOURCE", font=MONO, size=12, color=FAINT,
             align=PP_ALIGN.CENTER, tracking=8)
    # Chrome
    chrome(slide, 1)


def build_problem(slide):
    set_slide_bg(slide, WHITE)
    # Eyebrow
    add_text(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.4),
             "THE GAP", font=MONO, size=14, color=PRIMARY,
             align=PP_ALIGN.CENTER, tracking=16)
    # 2.6B huge number
    add_text(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(3.6),
             "2.6B", font=DISPLAY, size=300, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER, line_spacing=0.95)
    # Sub
    add_multi(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(0.7),
              [
                  {'text': 'people have ', 'size': 28, 'color': INK},
                  {'text': 'no internet', 'size': 28, 'color': INK, 'bold': True},
                  {'text': '.', 'size': 28, 'color': INK},
              ],
              align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(6.0), SLIDE_W, Inches(0.4),
             "300 million of them are students.",
             font=SANS, size=22, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(6.85), SLIDE_W, Inches(0.3),
             "SOURCE: UN ITU · UNESCO, 2024",
             font=MONO, size=11, color=FAINT, align=PP_ALIGN.CENTER, tracking=4)
    chrome(slide, 2)


def build_solution(slide):
    set_slide_bg(slide, WHITE)
    add_text(slide, Inches(0.7), Inches(0.9), SLIDE_W, Inches(0.4),
             "WHAT IT IS", font=MONO, size=14, color=PRIMARY,
             align=PP_ALIGN.LEFT, tracking=16)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(1.6),
             "A classroom that works\nwithout the internet.",
             font=DISPLAY, size=56, bold=True, color=INK,
             line_spacing=1.05)

    # Three pillars
    pillar_y = Inches(4.2)
    pillar_h = Inches(2.4)
    gap = Inches(0.3)
    total_w = SLIDE_W - Inches(1.4)
    pw = (total_w - gap * 2) / 3
    pillars = [
        ("Cached", "PWA service worker holds the whole app on the device."),
        ("Local",  "Lessons and quiz scores live in IndexedDB. No network needed."),
        ("Mesh",   "WebRTC syncs devices on the same WiFi, peer-to-peer."),
    ]
    for i, (title, sub) in enumerate(pillars):
        x = Inches(0.7) + (pw + gap) * i
        # card
        add_rect(slide, x, pillar_y, pw, pillar_h, SOFT, radius=Inches(0.18))
        # icon swatch
        icon_size = Inches(0.55)
        add_rect(slide, x + Inches(0.32), pillar_y + Inches(0.32),
                 icon_size, icon_size, PRIMARY, radius=Inches(0.12))
        # title
        add_text(slide, x + Inches(0.32), pillar_y + Inches(1.05),
                 pw - Inches(0.6), Inches(0.5),
                 title, font=DISPLAY, size=22, bold=True, color=INK)
        # sub
        add_text(slide, x + Inches(0.32), pillar_y + Inches(1.55),
                 pw - Inches(0.6), Inches(0.8),
                 sub, font=SANS, size=14, color=MUTED, line_spacing=1.35)

    chrome(slide, 3)


def build_cut(slide, line, rec_label):
    set_slide_bg(slide, BLACK)
    add_text(slide, Inches(0), Inches(2.8), SLIDE_W, Inches(1.2),
             line, font=DISPLAY, size=64, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    # rec dot + label
    rec_y = Inches(4.3)
    dot = add_circle(slide, Inches(5.3), rec_y + Inches(0.18),
                     Inches(0.1), BAD)
    add_text(slide, Inches(5.5), rec_y, Inches(3.5), Inches(0.4),
             rec_label, font=MONO, size=14, bold=True, color=BAD,
             align=PP_ALIGN.LEFT, tracking=10)
    chrome(slide, None, on_dark=True)


def build_mesh(slide):
    set_slide_bg(slide, WHITE)
    add_text(slide, Inches(0.7), Inches(0.9), SLIDE_W, Inches(0.4),
             "HOW THE SYNC SPREADS", font=MONO, size=14, color=PRIMARY,
             align=PP_ALIGN.LEFT, tracking=16)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(1.4),
             "One teacher publishes.\nEvery nearby device catches up in seconds.",
             font=DISPLAY, size=36, bold=True, color=INK,
             line_spacing=1.1)

    # Mesh diagram
    cx = Inches(6.67)
    cy = Inches(5.0)
    R = Inches(1.5)
    # Outer halo (very faint)
    add_circle(slide, cx, cy, R + Inches(0.5), WHITE, line=RULE, line_width=0.5)
    add_circle(slide, cx, cy, R, WHITE, line=RULE, line_width=0.5)
    # Spokes — angles in standard math (0=east, ccw +), Y goes DOWN in PPT
    # so a positive sin places the node BELOW the hub.
    import math
    students = []
    # Top-left, top-right, bottom-left, bottom-right, bottom-center
    angles = [210, -30, 150, 30, 90]
    labels = ['S1', 'S2', 'S3', 'S4', 'S5']
    for ang, lbl in zip(angles, labels):
        rad = math.radians(ang)
        sx = cx + int(R * math.cos(rad))
        sy = cy + int(R * math.sin(rad))
        add_line(slide, cx, cy, sx, sy, color=PRIMARY, weight=1.25)
        students.append((sx, sy, lbl))
    # Student nodes (over the lines)
    node_r = Inches(0.32)
    for sx, sy, lbl in students:
        c = add_circle(slide, sx, sy, node_r, WHITE, line=PRIMARY, line_width=2)
        # Label sits on top
        add_text(slide, sx - Inches(0.4), sy - Inches(0.15),
                 Inches(0.8), Inches(0.35),
                 lbl, font=SANS, size=11, bold=True, color=PRIMARY,
                 align=PP_ALIGN.CENTER)
    # Teacher hub
    hub_r = Inches(0.55)
    add_circle(slide, cx, cy, hub_r, PRIMARY)
    add_text(slide, cx - Inches(0.7), cy - Inches(0.16),
             Inches(1.4), Inches(0.35),
             "TEACHER", font=SANS, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # Body explanation (left column under headline)
    add_multi(slide, Inches(0.7), Inches(3.4), Inches(5.4), Inches(3.5),
              [
                  {'text': 'Heartbeat returns peers. New devices auto-connect. ', 'size': 16, 'color': MUTED},
                  {'text': 'When they connect, they swap peer lists.', 'size': 16, 'color': MUTED},
                  {'text': '\n\nA lesson posted on one tablet propagates to every other tablet on the local WiFi ', 'size': 16, 'color': MUTED},
                  {'text': 'within 30 seconds', 'size': 16, 'color': INK, 'bold': True},
                  {'text': '.', 'size': 16, 'color': MUTED},
              ],
              line_spacing=1.45)

    chrome(slide, 6)


def build_impact(slide):
    set_slide_bg(slide, WHITE)
    add_text(slide, Inches(0.7), Inches(0.9), SLIDE_W, Inches(0.4),
             "REACH + COST", font=MONO, size=14, color=PRIMARY,
             align=PP_ALIGN.LEFT, tracking=16)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(1.4),
             "Built on web standards.\nFree to ship.",
             font=DISPLAY, size=44, bold=True, color=INK,
             line_spacing=1.1)

    # 4 stat tiles
    tile_y = Inches(3.8)
    tile_h = Inches(1.6)
    gap = Inches(0.22)
    total_w = SLIDE_W - Inches(1.4)
    tw = (total_w - gap * 3) / 4
    stats = [
        ("$0",   "To deploy on Render's free tier"),
        ("$60",  "For a Raspberry Pi school server"),
        ("4",    "Translated UI languages"),
        ("0 KB", "Internet needed to take a quiz after install"),
    ]
    for i, (num, lbl) in enumerate(stats):
        x = Inches(0.7) + (tw + gap) * i
        add_rect(slide, x, tile_y, tw, tile_h, PANEL, line=RULE,
                 radius=Inches(0.16))
        add_text(slide, x + Inches(0.3), tile_y + Inches(0.25),
                 tw - Inches(0.6), Inches(0.7),
                 num, font=DISPLAY, size=40, bold=True, color=PRIMARY,
                 line_spacing=1.0)
        add_text(slide, x + Inches(0.3), tile_y + Inches(1.0),
                 tw - Inches(0.6), Inches(0.55),
                 lbl, font=SANS, size=12, color=MUTED, line_spacing=1.3)

    # SDG chips
    chip_y = Inches(5.85)
    chip_h = Inches(0.55)
    chips = [
        ("4",  "Quality education"),
        ("9",  "Infrastructure"),
        ("10", "Reduced inequalities"),
    ]
    x = Inches(0.7)
    for num, label in chips:
        # Estimate chip width from label length
        cw = Inches(0.9) + Inches(0.16) * len(label)
        # chip background pill
        add_rect(slide, x, chip_y, cw, chip_h, SOFT, radius=Inches(0.5))
        # number badge
        badge_r = Inches(0.22)
        bx = x + Inches(0.25)
        by = chip_y + chip_h / 2
        add_circle(slide, bx + badge_r, by, badge_r, PRIMARY)
        add_text(slide, bx, by - Inches(0.16), badge_r * 2, Inches(0.32),
                 num, font=DISPLAY, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        # label
        add_text(slide, x + Inches(0.78), chip_y + Inches(0.12),
                 cw - Inches(0.9), Inches(0.35),
                 label, font=SANS, size=14, bold=True, color=PRIMARY,
                 align=PP_ALIGN.LEFT)
        x = x + cw + Inches(0.2)

    chrome(slide, 7)


def build_closing(slide):
    set_slide_bg(slide, WHITE)
    # Logo mark
    mark_x = Inches(6.05)
    mark_y = Inches(1.7)
    add_rect(slide, mark_x, mark_y, Inches(0.7), Inches(0.7), PRIMARY,
             radius=Inches(0.18))
    add_circle(slide, mark_x + Inches(1.0), mark_y + Inches(0.35),
               Inches(0.32), WHITE, line=PRIMARY, line_width=2.4)

    add_text(slide, Inches(0), Inches(2.9), SLIDE_W, Inches(1.0),
             "Try it.", font=DISPLAY, size=56, bold=True, color=INK,
             align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0), Inches(4.2), SLIDE_W, Inches(0.6),
             "offlinefirst.onrender.com", font=MONO, size=22, bold=True,
             color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(4.95), SLIDE_W, Inches(0.6),
             "github.com/EshanthPen/offlinefirst", font=MONO, size=22,
             bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0), Inches(6.6), SLIDE_W, Inches(0.4),
             "OFFLINEFIRST  ·  YCS 2026", font=MONO, size=11, color=FAINT,
             align=PP_ALIGN.CENTER, tracking=8)
    chrome(slide, 8)


def chrome(slide, num, on_dark=False):
    """Top-right slide counter + top-left brand wordmark."""
    color = FAINT if not on_dark else RGBColor(0x9A, 0xA0, 0xA6)
    brand_color = PRIMARY if not on_dark else RGBColor(0xAE, 0xCB, 0xFA)
    add_text(slide, Inches(0.4), Inches(0.25), Inches(3), Inches(0.3),
             "offlinefirst", font=DISPLAY, size=11, bold=True,
             color=brand_color, align=PP_ALIGN.LEFT)
    if num is not None:
        add_text(slide, SLIDE_W - Inches(2.0), Inches(0.25),
                 Inches(1.6), Inches(0.3),
                 f"{num} / 8", font=MONO, size=11, color=color,
                 align=PP_ALIGN.RIGHT, tracking=8)
    else:
        add_text(slide, SLIDE_W - Inches(2.0), Inches(0.25),
                 Inches(1.6), Inches(0.3),
                 "CUT TO APP", font=MONO, size=11, color=BAD,
                 align=PP_ALIGN.RIGHT, tracking=10)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    # Slide 1 — Title
    s = prs.slides.add_slide(blank)
    build_title(s)
    add_notes(s,
              "Title — hold 3 beats before speaking.",
              "Silent. Three beats of music. Logo lands.",
              "Hold for 3 seconds before pressing →. No talking yet.")

    # Slide 2 — Problem
    s = prs.slides.add_slide(blank)
    build_problem(s)
    add_notes(s,
              "The problem.",
              '"Two-point-six billion people have no internet. Three hundred million of them are students."',
              'Pause after "no internet". Let the number sit on screen.')

    # Slide 3 — Solution
    s = prs.slides.add_slide(blank)
    build_solution(s)
    add_notes(s,
              "The solution.",
              '"OfflineFirst is a classroom that works without the internet. Three pieces of standard web tech: the app caches itself, lessons live in the browser database, devices on the same WiFi share over WebRTC."',
              "Don't overexplain. The demo is the proof. ~10 seconds.")

    # Slide 4 — Cut to demo 1
    s = prs.slides.add_slide(blank)
    build_cut(s, "Let's see it.", "REC  ·  ROLL APP DEMO")
    add_notes(s,
              "CUT to app — demo 1 (30s of screen recording).",
              '"Watch a student set up the app, open a lesson, and take a quiz."',
              "SWITCH to screen recording.\n  1. Onboarding (welcome → student → name 'Aminata' → grade)\n  2. Click Mathematics in sidebar\n  3. Open Introduction to Fractions\n  4. Scroll lesson · tap one Listen button\n  5. Click Start quiz → answer 2 questions → see result\nSWITCH back to deck.")

    # Slide 5 — Cut to demo 2
    s = prs.slides.add_slide(blank)
    build_cut(s, "Now turn off the WiFi.", "REC  ·  OFFLINE DEMO")
    add_notes(s,
              "CUT to app — offline demo (25s, MONEY SHOT).",
              '"Now I turn the internet completely off. The app still works. I take a quiz. When I come back online, the score appears on the teacher\'s dashboard."',
              "SWITCH to screen recording.\n  1. DevTools → Network → check Offline (show the throttling badge)\n  2. Reload the page → app still loads\n  3. Take any quiz, finish it, see the green score hero\n  4. DevTools → uncheck Offline\n  5. Wait 5s, switch to Teacher view (Settings → restart, or onboard as teacher in another tab)\n  6. Show the new score in Results\nDo not rush. This is the proof.")

    # Slide 6 — Mesh
    s = prs.slides.add_slide(blank)
    build_mesh(s)
    add_notes(s,
              "How the sync spreads.",
              '"When one teacher publishes a lesson, the server tells every other active device about it. They swap peer lists. New lessons reach every tablet on the local WiFi within thirty seconds."',
              "Pause on the diagram. Hands off keyboard while you talk.")

    # Slide 7 — Impact
    s = prs.slides.add_slide(blank)
    build_impact(s)
    add_notes(s,
              "Reach + cost.",
              '"Zero dollars to deploy on Render\'s free tier. Sixty dollars for a Raspberry Pi at a school. Four translated UI languages. Zero kilobytes of internet needed to take a quiz after first install. Three UN Sustainable Development Goals."',
              "Speak each number deliberately. Land on the SDGs.")

    # Slide 8 — Closing
    s = prs.slides.add_slide(blank)
    build_closing(s)
    add_notes(s,
              "Closing.",
              '"Live at offlinefirst-dot-onrender-dot-com. Open source. Built for the Young Coders Sphere competition."',
              "Hold for 3 seconds after speaking. Let the URL be readable.")

    out = Path(__file__).parent / "offlinefirst-pitch.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({out.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    main()
